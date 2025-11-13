import argparse
import json
from datetime import datetime
from pathlib import Path
from typing import Dict, Optional


def find_latest_op(artifacts_root: Path) -> Optional[Path]:
    ops = [p for p in artifacts_root.glob("*/") if p.is_dir()]
    if not ops:
        return None
    return sorted(ops, key=lambda p: p.stat().st_mtime, reverse=True)[0]


def load_metrics(op_dir: Path) -> Dict[str, Dict[str, object]]:
    metrics_path = op_dir / "output" / "report" / "metrics_summary.json"
    if metrics_path.exists():
        return json.loads(metrics_path.read_text(encoding="utf-8"))
    return {}


def load_manual_periods(op_dir: Path) -> Dict[str, list]:
    m = load_metrics(op_dir)
    # support both nested and flat
    data = m.get("metrics", m)
    y2025 = data.get("2025", {})
    return {"2025": y2025.get("manual_periods", [])}


def add_title_slide(prs, title_text: str, subtitle_text: str = ""):
    title_layout = prs.slide_layouts[0]  # Title
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title_text
    if subtitle_text:
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text


def add_metrics_table_slide(prs, metrics: Dict[str, Dict[str, object]]):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Key Metrics: 2022 vs 2025"
    # Build table data
    years = ["2022", "2025"]
    rows = [
        "total_applications",
        "total_interviews",
        "total_outreach",
        "applications_per_active_day",
        "median_days_between_applications",
        "median_days_to_next_interview",
        "interviews_per_application",
    ]
    n_rows = len(rows) + 1
    n_cols = 1 + len(years) + 1  # metric + years + delta
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    headers = ["Metric"] + years + ["Δ 2025-2022"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    def fmt(v):
        if isinstance(v, float):
            return f"{v:.2f}"
        return "" if v is None else str(v)
    for r, key in enumerate(rows, start=1):
        table.cell(r, 0).text = key
        v22 = metrics.get("2022", {}).get(key)
        v25 = metrics.get("2025", {}).get(key)
        table.cell(r, 1).text = fmt(v22)
        table.cell(r, 2).text = fmt(v25)
        delta = None
        try:
            if isinstance(v22, (int, float)) and isinstance(v25, (int, float)):
                delta = v25 - v22
        except Exception:
            pass
        table.cell(r, 3).text = fmt(delta)


def add_image_slide(prs, title: str, image_path: Path):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.5), width=Inches(9))


def add_two_images_slide(prs, title: str, left_image: Path, right_image: Path):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(left_image), Inches(0.5), Inches(1.5), width=Inches(4.25))
    slide.shapes.add_picture(str(right_image), Inches(5.25), Inches(1.5), width=Inches(4.25))


def add_bulleted_slide(prs, title: str, bullets: list[str]):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    tf = tx.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0


def add_two_column_text_slide(prs, title: str, left_title: str, left_points: list[str], right_title: str, right_points: list[str]):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    # Left column
    l_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.25), Inches(0.4)).text_frame
    l_title.text = left_title
    l_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.25), Inches(4.1)).text_frame
    l_box.clear()
    for i, t in enumerate(left_points):
        p = l_box.add_paragraph() if i > 0 else l_box.paragraphs[0]
        p.text = t
        p.level = 0
    # Right column
    r_title = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.25), Inches(0.4)).text_frame
    r_title.text = right_title
    r_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.9), Inches(4.25), Inches(4.1)).text_frame
    r_box.clear()
    for i, t in enumerate(right_points):
        p = r_box.add_paragraph() if i > 0 else r_box.paragraphs[0]
        p.text = t
        p.level = 0


def add_quote_slide(prs, title: str, quote: str):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(3)).text_frame
    box.text = f"“{quote}”"


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate comparative PowerPoint deck")
    parser.add_argument("--artifacts-root", default="artifacts/operations", help="Artifacts root")
    parser.add_argument("--op-dir", default=None, help="Specific operation dir")
    args = parser.parse_args()

    artifacts_root = Path(args.artifacts_root)
    op_dir = Path(args.op_dir) if args.op_dir else find_latest_op(artifacts_root)
    if not op_dir or not op_dir.exists():
        print("No operation directory found.")
        return

    charts_dir = op_dir / "output" / "charts"
    metrics = load_metrics(op_dir)
    manual_periods = load_manual_periods(op_dir)

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Title with provided headline
    add_title_slide(prs, "From Effort to Insight: Job Searching in the Age of AI", datetime.utcnow().strftime("Generated %Y-%m-%d %H:%M UTC"))

    # Opening Scene — 2022: The Manual Grind
    add_bulleted_slide(
        prs,
        "2022: The Manual Grind",
        [
            "Late nights: LinkedIn on one tab, motivation on another",
            "~100–160 applications, mostly manual",
            "Burnout curve (placeholder): add moodline if available",
            "Working hard, not necessarily smart",
        ],
    )

    # 2025: Automation Without Apathy (show cumulative apps vs interviews or inbound role)
    add_bulleted_slide(
        prs,
        "2025: Automation Without Apathy",
        [
            "1,346+ applications with lower emotional cost",
            "Spike in output, steadier mood (placeholder)",
            "Key: ‘I stopped writing applications; I designed systems that wrote for me.’",
            "Not doing less — doing what matters",
        ],
    )

    # Two-column reflection — The Human Operating System
    add_two_column_text_slide(
        prs,
        "The Human Operating System",
        "Psychology",
        [
            "Imposter syndrome evolves; it doesn’t vanish",
            "Efficiency ≠ Laziness — it’s strategic compassion for time",
            "Redefine ‘hard work’ for the cognitive age",
        ],
        "Visual",
        [
            "Placeholder: Two-axis diagram (Energy vs Meaning)",
            "Annotate inflection points (system redesigns, CV iterations)",
        ],
    )

    # Data speaks + Human speaks (metrics overview)
    add_metrics_table_slide(prs, metrics.get("metrics", metrics))
    add_bulleted_slide(
        prs,
        "Signals & Spikes",
        [
            "Apps per month (2022 vs 2025)",
            "Interview conversion rate (per year)",
            "Manual vs AI-assisted ratio (placeholder)",
            "Spikes map to system redesigns / CV iterations",
        ],
    )

    # Cumulative Apps vs Interviews per year (if images exist)
    for y in (2022, 2025):
        img = charts_dir / f"cumulative_apps_vs_interviews_{y}.png"
        if img.exists():
            add_image_slide(prs, f"Cumulative Applications vs Interviews - {y}", img)

    # Cumulative Apps with Inbound Spikes per year (side-by-side if both exist)
    img22 = charts_dir / "cumulative_apps_with_inbound_outreach_spikes_2022.png"
    img25 = charts_dir / "cumulative_apps_with_inbound_outreach_spikes_2025.png"
    if img22.exists() and img25.exists():
        add_two_images_slide(prs, "Cumulative Apps + Inbound Outreach Spikes", img22, img25)
    else:
        if img22.exists():
            add_image_slide(prs, "Cumulative Apps + Inbound Outreach Spikes - 2022", img22)
        if img25.exists():
            add_image_slide(prs, "Cumulative Apps + Inbound Outreach Spikes - 2025", img25)

    # Manual periods (2025): add one slide per period
    for idx, period in enumerate(manual_periods.get("2025", []), start=1):
        img = charts_dir / f"period_2025_{idx}_cumulative_apps_vs_inbound_role.png"
        if img.exists():
            title = period.get("name") or f"Period {idx}"
            add_image_slide(prs, f"{title} ({period.get('start_date')} to {period.get('end_date')})", img)

    # Manual vs AI-assisted
    mvai = charts_dir / "manual_vs_ai_assisted.png"
    if mvai.exists():
        add_image_slide(prs, "Manual vs AI-assisted Applications", mvai)

    # Facing the Critics (Q&A framing)
    add_two_column_text_slide(
        prs,
        "Facing the Critics",
        "Question",
        ["‘Aren’t you just letting AI do all the work?’"],
        "Reframe",
        [
            "‘No. I designed workflows so I could focus on thinking, connecting, learning.’",
            "‘AI took over repetition; I took over reflection.’",
        ],
    )

    # Closing
    add_bulleted_slide(
        prs,
        "The Future of Work — Personal Takeaways",
        [
            "It’s not who works hardest; it’s who learns fastest",
            "Authenticity remains the differentiator in a sea of automation",
            "‘I didn’t just change how I apply. I changed what effort means.’",
        ],
    )

    # Cumulative Apps vs Inbound Role Contacts per year
    for y in (2022, 2025):
        img = charts_dir / f"cumulative_apps_vs_inbound_role_{y}.png"
        if img.exists():
            add_image_slide(prs, f"Cumulative Apps vs Inbound Role Contacts - {y}", img)

    out_dir = op_dir / "output"
    out_path = out_dir / "Report_2022_vs_2025.pptx"
    prs.save(str(out_path))
    print(json.dumps({
        "operation_id": op_dir.name,
        "pptx": str(out_path)
    }, indent=2))


if __name__ == "__main__":
    # Lazy import requires Inches from pptx.util. Define here for module-level helpers above.
    from pptx.util import Inches
    main()


